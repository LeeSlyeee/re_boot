"""
Phase 2-2: 적응형 콘텐츠 분기 API Views
"""
import os
import openai
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.permissions import IsAuthenticated
from django.shortcuts import get_object_or_404
from django.utils import timezone

from .models import AdaptiveContent, LectureMaterial, LiveSession, PlacementResult


class GenerateAdaptiveView(APIView):
    """POST /api/learning/materials/{id}/generate-adaptive/ — 3레벨 변형 생성"""
    permission_classes = [IsAuthenticated]

    def post(self, request, pk):
        material = get_object_or_404(LectureMaterial, id=pk, uploaded_by=request.user)

        # 이미 생성되어 있으면 반환
        existing = AdaptiveContent.objects.filter(source_material=material)
        if existing.filter(status__in=['DRAFT', 'APPROVED']).count() >= 3:
            return Response({
                'message': '이미 3레벨 모두 생성되어 있습니다.',
                'levels': [{'id': a.id, 'level': a.level, 'status': a.status} for a in existing],
            })

        # 텍스트 추출
        source_text = self._extract_text(material)
        if not source_text:
            return Response({
                'error': f'텍스트를 추출할 수 없습니다. ({material.file_type} 형식)\n'
                         f'지원 형식: MD, TXT, PDF, PPT/PPTX, DOCX\n'
                         f'파일이 올바른 형식인지 확인해주세요. (.doc 구형식은 .docx로 변환 필요)'
            }, status=status.HTTP_400_BAD_REQUEST)

        # 3레벨 변형 생성 (대상별 맞춤 프롬프트)
        levels_config = {
            1: {
                'label': 'Level 1 - 쉽게 이해하기',
                'system': """당신은 **초등학생에게 설명하는 친절한 선생님**입니다.
아래 원본 교안의 **주제와 과목을 정확히 파악**하고, **초등학생도 이해할 수 있는 수준**으로 변형해주세요.

🎯 변형 원칙:
1. **쉬운 말만 사용**: 어려운 전문 용어는 모두 쉬운 말로 바꿔주세요.
2. **일상 비유**: 각 개념을 일상생활에서 볼 수 있는 것에 비유해서 설명하세요.
3. **쉬운 예시**: 초등학생이 공감할 수 있는 쉬운 예시를 많이 넣어주세요.
4. **이모지 활용**: 🍕🎒✏️ 등 이모지를 적극 활용해서 재미있게 구성하세요.
5. **짧은 문장**: 한 문장은 20자 이내로, 짧고 명확하게 쓰세요.
6. **핵심 정리**: 끝에 "📌 오늘 배운 것 정리!" 섹션을 추가하세요.

⛔ 금지 사항:
- **어려운 개념, 예외 사항, 주의할 점은 모두 빼세요.** 핵심만 쉽게 설명하세요.
- 복잡한 규칙이나 심화 내용은 절대 넣지 마세요.
- **원본 교안의 과목/주제를 벗어나는 내용은 절대 넣지 마세요.**

마크다운 형식으로 작성하세요.""",
            },
            2: {
                'label': 'Level 2 - 핵심 정리',
                'system': """당신은 **중학생을 가르치는 학교 선생님**입니다.
아래 원본 교안의 **주제와 과목을 정확히 파악**하고, **중학생(간단한 기초 교육을 받은 학생)** 수준으로 변형해주세요.

🎯 변형 원칙:
1. **원본 구조 유지**: 원본의 목차와 흐름을 유지하세요.
2. **개념 설명 보강**: 각 개념을 "왜 이런 규칙이 있는지" 간단한 이유와 함께 설명하세요.
3. **적절한 예시**: 각 개념마다 이해를 돕는 예시를 2~3개 추가하세요.
4. **핵심 강조**: 중요한 부분은 **볼드**로 표시하세요.
5. **연습 문제**: "📝 연습 문제" 섹션에 쉬운~보통 난이도의 문제 3~5개를 추가하세요.
   - 각 문제에 정답과 간단한 해설을 포함하세요.
6. **💡 팁**: 자주 틀리는 부분을 "💡 이것만 기억하세요" 박스로 정리하세요.

⚠️ 중요:
- 요약이 아니라 **보강**입니다. 원본보다 길게 작성하세요.
- **원본 교안의 과목/주제를 벗어나는 내용은 절대 넣지 마세요.**

마크다운 형식으로 작성하세요.""",
            },
            3: {
                'label': 'Level 3 - 심화 완성',
                'system': """당신은 **성인 학습자를 위한 전문 강사**입니다.
아래 원본 교안의 **주제와 과목을 정확히 파악**하고, **기본적인 교육을 받은 성인 학습자** 수준으로 변형해주세요.

🎯 변형 원칙:
1. **체계적 정리**: 원본 내용을 논리적이고 체계적으로 정리하세요.
2. **배경 지식 추가**: 각 개념의 "왜(Why)"와 원리를 간결하게 설명하세요.
3. **실전 활용**: 실제 상황에서 어떻게 활용하는지 실용적 사례를 제시하세요.
4. **주의사항**: 흔히 틀리는 부분과 예외 사항을 "⚠️ 주의" 박스로 정리하세요.
5. **심화 연습 문제**: "📝 실전 문제" 섹션에 응용 문제 5개를 추가하세요.
   - 각 문제에 정답과 상세한 해설을 포함하세요.
6. **시험 대비**: "🎯 시험에 나올 수 있는 포인트"를 정리하세요.
7. **추가 학습**: "📚 더 알아보기" 섹션에 관련 참고 자료를 추천하세요.

⚠️ 중요:
- 요약이 아니라 **확장**입니다. 원본보다 **2배 이상** 길고 깊이 있게 작성하세요.
- **원본 교안의 과목/주제를 벗어나는 내용은 절대 넣지 마세요.**

마크다운 형식으로 작성하세요.""",
            },
        }

        results = []
        client = openai.OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

        for level, config in levels_config.items():
            # 중복 방지
            if AdaptiveContent.objects.filter(source_material=material, level=level).exists():
                ac = AdaptiveContent.objects.get(source_material=material, level=level)
                results.append({'id': ac.id, 'level': level, 'status': ac.status})
                continue

            ac = AdaptiveContent.objects.create(
                source_material=material,
                level=level,
                title=f"{material.title} ({config['label']})",
                status='GENERATING',
            )

            try:
                response = client.chat.completions.create(
                    model='gpt-4o',
                    messages=[
                        {'role': 'system', 'content': config['system']},
                        {'role': 'user', 'content': f'원본 교안:\n\n{source_text[:6000]}'},
                    ],
                    temperature=0.6,
                    max_tokens=4000,
                )

                ac.content = response.choices[0].message.content.strip()
                ac.status = 'DRAFT'
                ac.save()
                results.append({'id': ac.id, 'level': level, 'status': 'DRAFT'})

            except Exception as e:
                ac.status = 'DRAFT'
                ac.content = f'[생성 실패: {str(e)}]'
                ac.save()
                results.append({'id': ac.id, 'level': level, 'status': 'DRAFT', 'error': str(e)})

        return Response({'levels': results}, status=status.HTTP_201_CREATED)

    def _extract_text(self, material):
        """교안 텍스트 추출 (1차: MD만 지원)"""
        if material.file_type == 'MD' and material.file:
            try:
                material.file.open('r')
                text = material.file.read()
                material.file.close()
                if isinstance(text, bytes):
                    text = text.decode('utf-8')
                return text
            except Exception:
                pass

        # PDF: PyMuPDF 시도
        if material.file_type == 'PDF' and material.file:
            try:
                import fitz
                doc = fitz.open(material.file.path)
                text = ''
                for page in doc:
                    text += page.get_text()
                doc.close()
                return text if text.strip() else None
            except ImportError:
                pass

        # PPT: python-pptx 시도
        if material.file_type == 'PPT' and material.file:
            try:
                from pptx import Presentation
                prs = Presentation(material.file.path)
                text = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text + '\n'
                return text if text.strip() else None
            except ImportError:
                pass

        # DOCX: python-docx 시도
        if material.file_type == 'DOCX' and material.file:
            try:
                from docx import Document
                import io
                # Django FileField를 통해 바이너리로 읽어서 전달 (한글 경로 이슈 방지)
                material.file.open('rb')
                file_bytes = material.file.read()
                material.file.close()
                doc = Document(io.BytesIO(file_bytes))
                text = ''
                for para in doc.paragraphs:
                    text += para.text + '\n'
                # 테이블 내용도 추출
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text += cell.text + ' '
                        text += '\n'
                return text if text.strip() else None
            except Exception as e:
                print(f"⚠️ DOCX 텍스트 추출 실패: {e}")
                # .doc (구 바이너리 형식)일 수 있음 → 안내 메시지
                return None

        # TXT: 직접 읽기
        if material.file_type == 'TXT' and material.file:
            try:
                material.file.open('rb')
                raw = material.file.read()
                material.file.close()
                # 인코딩 자동 감지
                for enc in ['utf-8', 'euc-kr', 'cp949', 'latin-1']:
                    try:
                        text = raw.decode(enc)
                        return text
                    except UnicodeDecodeError:
                        continue
                return raw.decode('utf-8', errors='replace')
            except Exception:
                pass

        return None


class ListAdaptiveView(APIView):
    """GET /api/learning/materials/{id}/adaptive/ — 교안의 레벨별 변형 목록"""
    permission_classes = [IsAuthenticated]

    def get(self, request, pk):
        material = get_object_or_404(LectureMaterial, id=pk)
        contents = AdaptiveContent.objects.filter(source_material=material)

        data = [
            {
                'id': ac.id,
                'level': ac.level,
                'level_label': ac.get_level_display(),
                'title': ac.title,
                'content': ac.content,
                'content_preview': ac.content[:200] + '...' if len(ac.content) > 200 else ac.content,
                'status': ac.status,
                'created_at': ac.created_at,
            }
            for ac in contents
        ]
        return Response({'adaptive_contents': data})


class ApproveAdaptiveView(APIView):
    """POST /api/learning/adaptive/{id}/approve/ — 변형 승인"""
    permission_classes = [IsAuthenticated]

    def post(self, request, pk):
        ac = get_object_or_404(AdaptiveContent, id=pk, source_material__uploaded_by=request.user)
        ac.status = 'APPROVED'
        ac.approved_at = timezone.now()
        ac.save()
        return Response({'ok': True, 'status': 'APPROVED'})


class MyContentView(APIView):
    """GET /api/learning/live/{id}/my-content/ — 학생 레벨에 맞는 자료 조회"""
    permission_classes = [IsAuthenticated]

    def get(self, request, pk):
        session = get_object_or_404(LiveSession, id=pk)

        # 학생 레벨 결정 (PlacementResult 기반)
        placement = PlacementResult.objects.filter(student=request.user).order_by('-created_at').first()
        level_map = {'BEGINNER': 1, 'INTERMEDIATE': 2, 'ADVANCED': 3}
        student_level = level_map.get(placement.level, 2) if placement else 2

        # 세션에 연결된 교안의 적응형 콘텐츠 조회
        materials = session.lecture.materials.all() if session.lecture else []
        my_contents = []

        for material in materials:
            # 내 레벨 콘텐츠 (승인된 것만)
            adaptive = AdaptiveContent.objects.filter(
                source_material=material,
                level=student_level,
                status='APPROVED',
            ).first()

            my_contents.append({
                'material_id': material.id,
                'material_title': material.title,
                'file_type': material.file_type,
                'adaptive_content': {
                    'id': adaptive.id,
                    'level': adaptive.level,
                    'level_label': adaptive.get_level_display(),
                    'title': adaptive.title,
                    'content': adaptive.content,
                } if adaptive else None,
                'original_url': material.file.url if material.file else '',
            })

        return Response({
            'student_level': student_level,
            'contents': my_contents,
        })
